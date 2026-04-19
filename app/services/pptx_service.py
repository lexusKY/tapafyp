from pptx import Presentation


def extract_text_from_pptx(pptx_path):
    try:
        prs = Presentation(pptx_path)
        text_parts = []

        for slide_number, slide in enumerate(prs.slides, start=1):
            slide_text = [f"--- Slide {slide_number} ---"]

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())

            text_parts.append("\n".join(slide_text))

        extracted_text = "\n\n".join(text_parts).strip()

        if not extracted_text:
            return None

        return extracted_text

    except Exception as e:
        print(f"PPTX extraction error: {e}")
        return None